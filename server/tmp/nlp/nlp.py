#必要なライブラリ
import pptx 
import re
import termextract.janome 
import termextract.core
from janome.tokenizer import Tokenizer
import collections
from sklearn.feature_extraction.text import TfidfVectorizer
import pandas as pd
import urllib.request
import os.path
import glob
from sklearn import model_selection
import joblib
import zipfile

#必要に応じて
#!pip install mecab-python3
#!pip install unidic-lite
import MeCab

#関数

#デフォルトの文字を省く
def sub(text): 
    text =re.sub('タイトル :', '', text)
    text =re.sub('サブ', '', text)
    text =re.sub('本文1 :', '', text)
    text =re.sub('本文2 :', '', text)
    text =re.sub('本文3 :', '', text)
    text =re.sub('図1 :', '', text)
    text =re.sub('図2 :', '', text)
    text =re.sub('図3 :', '', text)
    
    return text

####pptx(fname)の中身をテキスト抽出(sample.txtとして出力)
#形態素解析用に文ごとの配列を返す
def pptx_text(fname):
    f = open('sample.txt', 'w')
    #fname='sampleFile.pptx'
    txt = []
    prs = pptx.Presentation(fname)

    for i, sld in enumerate(prs.slides, start=1):

        #print(f'-- Page {i} --')

        for shp in sld.shapes:

            if shp.has_text_frame:
                #print (shp.text)
                shp.text = sub(shp.text)
                txt.append(shp.text)
                f.write(shp.text)

            if shp.has_table:
                tbl = shp.table
                row_count = len(tbl.rows)
                col_count = len(tbl.columns)
                for r in range(0, row_count):                 
                    text=''
                    for c in range(0, col_count):
                        cell = tbl.cell(r,c)
                        paragraphs = cell.text_frame.paragraphs 
                        for paragraph in paragraphs:
                            for run in paragraph.runs:
                                text+=run.text
                            text+=', '
                    print (text)
                    f.write(text)
            #print ()
    f.close()
    sentences = [s for s in txt if s != ' ']
    sentences = [s for s in sentences if s != '']
    return sentences
    
#sample.txtの中から重要度が高い順にした配列を返す
def keyword():
    # ファイルパス
    file_path = "sample.txt"
    
    # 日本語テキストの読み込み
    f = open(file_path, "r", encoding="utf-8")
    text = f.read()
    f.close
    words=[]
    # 形態素解析器で日本語処理
    t = Tokenizer()
    tokenize_text = t.tokenize(text)
    
    # Frequency生成＝複合語抽出処理（ディクショナリとリストの両方可)
    frequency = termextract.janome.cmp_noun_dict(tokenize_text)
    
    # FrequencyからLRを生成する
    lr = termextract.core.score_lr(
        frequency,
        ignore_words=termextract.mecab.IGNORE_WORDS,
        lr_mode=1, average_rate=1)
    
    # FrequencyとLRを組み合わせFLRの重要度を出す
    term_imp = termextract.core.term_importance(frequency, lr)
    
    # collectionsを使って重要度が高い順に表示
    data_collection = collections.Counter(term_imp)
    for cmp_noun, value in data_collection.most_common():
        #print(termextract.core.modify_agglutinative_lang(cmp_noun), value, sep="\t")
        words.append(termextract.core.modify_agglutinative_lang(cmp_noun))
    return words


#形態素解析をした後、形容詞と名詞を格納した配列を返す
def morpheme(sentences):
    # Tokenizerインスタンスの生成 
    t = Tokenizer()
    #print(text)
    # テキストを引数として、形態素解析の結果、名詞・形容詞(原形)のみを配列で抽出する関数を定義 
    def extract_words(text):
        tokens = t.tokenize(text)
        return [token.base_form for token in tokens 
            if token.part_of_speech.split(',')[0] in['名詞','形容詞']] 

    # それぞれの文章を単語リストに変換
    word_list = [extract_words(s) for s in sentences]
    words=[]
    for word in word_list:
        #print(word)
        words.append(word)
    return words

#TF-IDF分析に使う
def extract(text):
    words = []
    # Taggerオブジェクトを生成
    tokenizer = MeCab.Tagger("")
    tokenizer.parse("")
    # 単語の特徴リストを生成
    node = tokenizer.parseToNode(text)

    while node:
        # 品詞情報(node.feature)が名詞か形容詞ならば
        if node.feature.split(",")[0] == u"名詞" or node.feature.split(",")[0] == u"形容詞":
            # 単語(node.surface)をwordsに追加
            words.append(node.surface)
        node = node.next

    # 半角スペース区切りで文字列を結合
    text_result = ' '.join(words)
    return text_result

#TF-IDF分析(上位10単語を返す)
def tf_idf(): #filename...モデルの保存ファイルが含まれるzipファイル
    for filename in glob.glob('*.zip'):
        #zipファイル解凍
        zip_f = zipfile.ZipFile(filename)
        zip_f.extractall('./data')
        zip_f.close()
    # データフレームに表現
    X =joblib.load('./data/X.sav')
    # データフレームに表現
    values = X.toarray()
    feature_names = joblib.load('./data/finalized_model.sav')
    #feature_names = vectorizer.get_feature_names()
    df = pd.DataFrame(values, columns = feature_names,
                    index=["サンプル","心理学", "情報科学", "経済","社会学","医学","化学","数学","生物"])
    # 行列を転置
    df_0 = df[0:1].T
    # 値で降順ソート
    df_0 = df_0.sort_values(by="サンプル", ascending=False)
    return list(df_0.head(10).index)

if __name__ == '__main__':
    #一時的にtest.pptxにしてるので必要に応じて変更
    sentences = pptx_text('../pptx/data/test.pptx') #形態素解析に使用する配列 TF-IDF分析をするときにも必要なファイルを作るので必須
    key_words = keyword() #重要な順に単語を格納した配列
    words = morpheme(sentences)#形容詞と名詞を格納した配列
    df = tf_idf() #TF-IDF分析　上位10単語のリストを返す
    #print(df)
    