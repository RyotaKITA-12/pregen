import colorsys
import csv
import subprocess

import gensim
import MeCab
import pandas as pd
import pptx
from pptx.dml.color import RGBColor
from pptx.util import Cm
import spacy


def main(filename, colors, logo):
    path = '../data/'+filename+'.pptx'
    corrected_colors = correct_color(colors)

    read_ppt = pptx.Presentation(path)
    write_ppt = pptx.Presentation('../data/template.pptx')
    for i, slide in enumerate(read_ppt.slides):
        title, conclude, points = extract_items(path, slide)
        overwrite_template(write_ppt, i, logo,
                           corrected_colors, title, conclude, points, filename)


def extract_items(path, slide):
    points = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            if shape.text[:4] == "タイトル":
                title = shape.text[5:-1]
            if shape.text[:2] == "結論":
                conclude = shape.text[3:-1]
            if shape.text[:5] == "箇条書き1":
                points.append(shape.text[6:-1])
            if shape.text[:5] == "箇条書き2":
                points.append(shape.text[6:-1])
            if shape.text[:5] == "箇条書き3":
                points.append(shape.text[6:-1])

    return title, conclude, points


def overwrite_template(ppt, page, logo, colors, title, conclude, points, filename):

    slide = ppt.slides[page]
    slide.shapes[0].fill.solid()
    slide.shapes[0].fill.fore_color.rgb = RGBColor(*colors[0])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(*colors[1])
    #
    slide.shapes[1].text = title
    slide.shapes[1].text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(
        255, 255, 255)

    keyword = str(parsing_dependencies(conclude))
    text = conclude.rsplit(keyword, 1)
    slide.shapes[2].text = text[0]
    slide.shapes[2].text_frame.paragraphs[0].alignment = pptx.enum.text.PP_ALIGN.CENTER
    slide.shapes[2].text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(
        *colors[2])
    run = slide.shapes[2].text_frame.paragraphs[0].add_run()
    run.text = keyword
    run.font.color.rgb = RGBColor(*colors[0])
    run = slide.shapes[2].text_frame.paragraphs[0].add_run()
    run.text = text[1]
    run.font.color.rgb = RGBColor(*colors[2])
    slide.shapes[3].text_frame.paragraphs[0].text = points[0]
    slide.shapes[3].text_frame.paragraphs[1].text = points[1]
    slide.shapes[3].text_frame.paragraphs[2].text = points[2]
    run = slide.shapes[2].text_frame.paragraphs[0].add_run()
    slide.shapes[3].text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(
        *colors[2])
    slide.shapes[3].text_frame.paragraphs[1].runs[0].font.color.rgb = RGBColor(
        *colors[2])
    slide.shapes[3].text_frame.paragraphs[2].runs[0].font.color.rgb = RGBColor(
        *colors[2])
    slide.shapes.element.remove(slide.shapes[4].element)

    convert_pptx_to_zip(filename)

    slide.shapes.add_picture(
        '../data/unzipped/ppt/media/image1.png', Cm(4), Cm(8), Cm(9))

    slide.shapes.add_picture(
        '../data/img/'+logo, Cm(0.5), Cm(0.5), Cm(2))

    ppt.save("output.pptx")


def correct_color(colors):
    hsv_1 = list(colorsys.rgb_to_hsv(*colors[0]))
    corrected_hsv_1 = [hsv_1[0], 0.6, 230]
    colors[0] = list(map(int, colorsys.hsv_to_rgb(*corrected_hsv_1)))
    hsv_2 = list(colorsys.rgb_to_hsv(*colors[1]))
    corrected_hsv_2 = [hsv_2[0], hsv_2[1], 400]
    colors[1] = map(int, colorsys.hsv_to_rgb(*corrected_hsv_2))
    colors[1] = [col if col <= 255 else 255 for col in colors[1]]
    colors[1] = [col if 215 <= col else 215 for col in colors[1]]
    hsv_3 = list(colorsys.rgb_to_hsv(*colors[2]))
    corrected_hsv_3 = [hsv_3[0], hsv_3[1], 45]
    colors[2] = map(int, colorsys.hsv_to_rgb(*corrected_hsv_3))
    colors[2] = [col if 0 <= col else 0 for col in colors[2]]

    return colors


def convert_pptx_to_zip(filename):
    try:
        _ = subprocess.run(["rm", "-r", "-f", "../data/unzipped"])
    except FileNotFoundError:
        pass
    _ = subprocess.run(["mkdir", "../data/unzipped"])
    _ = subprocess.run(
        ["cp", "../data/"+filename+".pptx", "../data/"+filename+".zip"])
    _ = subprocess.run(
        ["unzip", "../data/"+filename+".zip", "-d", "../data/unzipped/"])


def parsing_dependencies(conclude):
    sp = spacy.load('ja_ginza')
    doc = sp(conclude)
    word_dic = {}
    for token in doc:
        word_dic[token] = 0
    for token in doc:
        word_dic[token.head] += 1
    max_k = max(word_dic, key=word_dic.get)
    return max_k


def choice_color_direct(input_word, color_patterns, colors):
    impressions = color_patterns['name']
    indexes = [i for i, name in enumerate(impressions) if name == input_word]
    selected_impressions = color_patterns.iloc[indexes]
    selected_names = list(selected_impressions['name'])

    result_colors = get_color_images(
        colors, selected_impressions, indexes, selected_names)

    return result_colors[0]


def get_color_images(colors, impressions, indexes, names):
    result_colors = []
    for i, _ in enumerate(indexes):
        used_colors = list(impressions.iloc[i][impressions.iloc[i] == 1].index)
        result_a_color = []
        for j, color in enumerate(used_colors):
            name_rgb = colors[colors['Hue/Tone'] == color]
            rgb = [int(name_rgb['R']), int(name_rgb['G']), int(name_rgb['B'])]
            result_a_color.append(rgb)
        result_colors.append(result_a_color)

    return result_colors


def extract_origin(word):
    m = MeCab.Tagger('-Ochasen')
    node = m.parseToNode(word)
    while node:
        origin = node.feature.split(",")[6]
        if origin == "*":
            node = node.next
        else:
            return origin


def create_similar(word):
    word_origin = extract_origin(word)
    impressions = list(dict.fromkeys(list(color_patterns['name'])))
    similar_dict = {}
    for impression in impressions:
        impression_origin = extract_origin(impression)
        try:
            similar_dict[impression] = model.wv.similarity(word_origin,
                                                           impression_origin)
        except KeyError:
            pass
    return sorted(similar_dict.items(), key=lambda x: x[1], reverse=True)


if __name__ == '__main__':
    colors = pd.read_csv('../data/color.csv')
    color_patterns = pd.read_csv('../data/color_img_scale.csv')
    color_patterns.rename(
        columns={'Unnamed: 0': 'name'},
        inplace=True)
    model = gensim.models.Word2Vec.load('../model/word2vec.gensim.model')
    word = input("word :")
    sim_dic = create_similar(word)
    colors = choice_color_direct(sim_dic[0][0], color_patterns, colors)
    filename = input("filename :")
    icons = {}
    with open('../data/icon.csv', mode='r') as file:
        reader = csv.reader(file)
        icons = {rows[1]: rows[0] for rows in reader}
    similar_dict = {}
    for icon in icons:
        try:
            similar_dict[icon] = model.wv.similarity(word, icon)
        except KeyError:
            pass
    logo = icons[sorted(similar_dict.items(),
                        key=lambda x: x[1], reverse=True)[0][0]]
    print(logo)
    main(filename, colors, logo)
