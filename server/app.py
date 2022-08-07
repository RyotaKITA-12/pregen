import collections
import colorsys
import csv
import glob
import joblib
import os
import subprocess
import zipfile

import gensim
from janome.tokenizer import Tokenizer
import MeCab
import pandas as pd
import pptx
from pptx.dml.color import RGBColor
from pptx.util import Cm
from pptx.util import Pt
import spacy
import termextract.janome
import termextract.core


def main(in_path, out_path):
    filename = os.path.splitext(os.path.basename(in_path))[0]
    sentences = pptx_text(in_path)
    keywords = keyword()
    words = morpheme(sentences)
    df = tf_idf()
    model = gensim.models.Word2Vec.load('./model/word2vec.gensim.model')
    df = ["猫", "自然", "二十歳", "大学", "ベッド"]

    icons = {}
    with open('./data/icon.csv', mode='r') as file:
        reader = csv.reader(file)
        icons = {rows[1]: rows[0] for rows in reader}
    similar_icon_dict = {}
    similar_word_dict = {}
    for word in df:
        for icon in icons:
            try:
                similar_icon_dict[icon] = model.wv.similarity(word, icon)
            except KeyError:
                pass
        similar_word_dict[word] = sorted(similar_icon_dict.items(),
                                         key=lambda x: x[1], reverse=True)[0][1]
    word = list(similar_word_dict)[0][0]
    for icon in icons:
        try:
            print(icon, model.wv.similarity(word, icon))
            similar_icon_dict[icon] = model.wv.similarity(word, icon)
        except KeyError:
            pass
    logo = icons[sorted(similar_icon_dict.items(),
                        key=lambda x: x[1], reverse=True)[0][0]]

    colors = pd.read_csv('../data/color.csv')
    color_patterns = pd.read_csv('../data/color_img_scale.csv')
    color_patterns.rename(
        columns={'Unnamed: 0': 'name'},
        inplace=True)
    sim_dic = create_similar(word)
    colors = choice_color_direct(sim_dic[0][0], color_patterns, colors)
    path = '../data/'+filename+'.pptx'
    corrected_colors = correct_color(colors)

    read_ppt = pptx.Presentation(path)

    for i, slide in enumerate(read_ppt.slides):
        slide_type, items = extract_items(path, slide)
        if slide_type["is_title"]:
            title, subtitle, affiliation, name, date = items
            overwrite_title(logo, corrected_colors, title,
                            subtitle, affiliation, name, date, filename, out_path)
        if slide_type["is_points"]:
            title, conclude, points = items
            overwrite_points(logo, corrected_colors, title,
                             conclude, points, filename, out_path)


def extract_items(path, slide):
    points = []
    slide_type = {"is_title": False, "is_points": False}
    subtitle, affiliation, name, date, conclude = "", "", "", "", ""
    for shape in slide.shapes:
        if shape.has_text_frame:
            print(shape.text)
            if shape.text[:4] == "タイトル":
                title = shape.text[5:-1]
            if shape.text[:6] == "サブタイトル":
                subtitle = shape.text[7:-1]
            if shape.text[:2] == "所属":
                affiliation = shape.text[3:-1]
            if shape.text[:2] == "結論":
                conclude = shape.text[3:-1]
            if shape.text[:2] == "日付":
                date = shape.text[3:-1]
            if shape.text[:2] == "名前":
                name = shape.text[3:-1]
            if shape.text[:5] == "箇条書き1":
                points.append(shape.text[6:-1])
            if shape.text[:5] == "箇条書き2":
                points.append(shape.text[6:-1])
            if shape.text[:5] == "箇条書き3":
                points.append(shape.text[6:-1])
            if shape.text[:5] == "箇条書き4":
                points.append(shape.text[6:-1])
            if shape.text[:5] == "箇条書き5":
                points.append(shape.text[6:-1])
            if shape.text[:5] == "箇条書き6":
                points.append(shape.text[6:-1])
            if shape.text[:5] == "箇条書き7":
                points.append(shape.text[6:-1])
            if shape.text[:5] == "箇条書き8":
                points.append(shape.text[6:-1])
            if shape.text[:5] == "箇条書き9":
                points.append(shape.text[6:-1])
            if shape.text[:4] == "大項目1":
                title = shape.text[5:-1]
            if shape.text[:4] == "大項目2":
                title = shape.text[5:-1]
            if shape.text == 'file:title':
                slide_type["is_title"] = True
            if shape.text == 'file:points':
                slide_type["is_points"] = True
    if slide_type["is_points"]:
        return slide_type, (title, conclude, points)
    if slide_type["is_title"]:
        return slide_type, (title, subtitle, affiliation, name, date)
    return slide_type, None


def overwrite_points(logo, colors, title, conclude, points, filename, out_path):
    ppt = pptx.Presentation('../templates/points.pptx')
    slide = ppt.slides[0]
    slide.shapes[0].fill.solid()
    slide.shapes[0].fill.fore_color.rgb = RGBColor(*colors[0])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(*colors[1])

    slide.shapes[1].text = title
    slide.shapes[1].text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(
        255, 255, 255)

    keyword = str(parsing_dependencies(conclude))
    text = conclude.rsplit(keyword, 1)
    slide.shapes[2].text = text[0]
    slide.shapes[2].text_frame.paragraphs[0].alignment = pptx.enum.text.PP_ALIGN.CENTER
    slide.shapes[2].text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(
        *colors[2])
    run = slide.shapes[2].text_frame.paragraphs[0].add_run()
    run.text = keyword
    run.font.color.rgb = RGBColor(*colors[0])
    run = slide.shapes[2].text_frame.paragraphs[0].add_run()
    run.text = text[1]
    run.font.color.rgb = RGBColor(*colors[2])
    slide.shapes[3].text_frame.paragraphs[0].text = points[0]
    slide.shapes[3].text_frame.paragraphs[1].text = points[1]
    slide.shapes[3].text_frame.paragraphs[2].text = points[2]
    run = slide.shapes[2].text_frame.paragraphs[0].add_run()
    slide.shapes[3].text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(
        *colors[2])
    slide.shapes[3].text_frame.paragraphs[1].runs[0].font.color.rgb = RGBColor(
        *colors[2])
    slide.shapes[3].text_frame.paragraphs[2].runs[0].font.color.rgb = RGBColor(
        *colors[2])
    slide.shapes.element.remove(slide.shapes[4].element)

    convert_pptx_to_zip(filename)

    slide.shapes.add_picture(
        '../data/unzipped/ppt/media/image1.png', Cm(4), Cm(8), Cm(9))

    slide.shapes.add_picture(
        '../data/img/'+logo, Cm(0.5), Cm(0.5), Cm(2))

    ppt.save(out_path)


def overwrite_title(logo, colors, title, subtitle, affiliation, name, date,
                    filename, out_path):
    ppt = pptx.Presentation('../templates/title.pptx')
    slide = ppt.slides[0]
    hsv_1 = list(colorsys.rgb_to_hsv(*colors[0]))
    hsv_2 = list(colorsys.rgb_to_hsv(*colors[1]))
    line_hsv = [hsv_1[0], 0.6, 150]
    title_hsv = [hsv_1[0], 0.9, 50]
    font_hsv = [hsv_2[0], 0.6, 50]
    line_color = list(map(int, colorsys.hsv_to_rgb(*line_hsv)))
    title_color = list(map(int, colorsys.hsv_to_rgb(*title_hsv)))
    font_color = list(map(int, colorsys.hsv_to_rgb(*font_hsv)))
    slide.shapes[2].text = title
    slide.shapes[2].text_frame.paragraphs[0].font.size = Pt(54)
    slide.shapes[2].text_frame.paragraphs[0].alignment = pptx.enum.text.PP_ALIGN.CENTER
    slide.shapes[2].text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(
        *title_color)
    slide.shapes[0].text = subtitle
    slide.shapes[0].text_frame.paragraphs[0].alignment = pptx.enum.text.PP_ALIGN.CENTER
    slide.shapes[0].text_frame.paragraphs[0].font.size = Pt(36)
    slide.shapes[0].text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(
        *colors[0])
    slide.shapes[1].line.color.rgb = RGBColor(*line_color)
    slide.shapes[3].text = affiliation
    slide.shapes[3].text_frame.paragraphs[0].alignment = pptx.enum.text.PP_ALIGN.CENTER
    slide.shapes[3].text_frame.paragraphs[0].font.size = Pt(24)
    slide.shapes[3].text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(
        *font_color)
    slide.shapes[4].text = name
    slide.shapes[4].text_frame.paragraphs[0].alignment = pptx.enum.text.PP_ALIGN.CENTER
    slide.shapes[4].text_frame.paragraphs[0].font.size = Pt(16)
    slide.shapes[4].text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(
        *font_color)
    slide.shapes[5].text = date
    slide.shapes[5].text_frame.paragraphs[0].alignment = pptx.enum.text.PP_ALIGN.RIGHT
    slide.shapes[5].text_frame.paragraphs[0].font.size = Pt(16)
    slide.shapes[5].text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(
        *font_color)

    ppt.save(out_path)


def correct_color(colors):
    hsv_1 = list(colorsys.rgb_to_hsv(*colors[0]))
    corrected_hsv_1 = [hsv_1[0], 0.6, 230]
    colors[0] = list(map(int, colorsys.hsv_to_rgb(*corrected_hsv_1)))
    hsv_2 = list(colorsys.rgb_to_hsv(*colors[1]))
    corrected_hsv_2 = [hsv_2[0], hsv_2[1], 400]
    colors[1] = map(int, colorsys.hsv_to_rgb(*corrected_hsv_2))
    colors[1] = [col if col <= 255 else 255 for col in colors[1]]
    colors[1] = [col if 215 <= col else 215 for col in colors[1]]
    hsv_3 = list(colorsys.rgb_to_hsv(*colors[2]))
    corrected_hsv_3 = [hsv_3[0], hsv_3[1], 45]
    colors[2] = map(int, colorsys.hsv_to_rgb(*corrected_hsv_3))
    colors[2] = [col if 0 <= col else 0 for col in colors[2]]

    return colors


def convert_pptx_to_zip(filename):
    try:
        _ = subprocess.run(["rm", "-r", "-f", "../data/unzipped"])
    except FileNotFoundError:
        pass
    _ = subprocess.run(["mkdir", "../data/unzipped"])
    _ = subprocess.run(
        ["cp", "../data/"+filename+".pptx", "../data/"+filename+".zip"])
    _ = subprocess.run(
        ["unzip", "../data/"+filename+".zip", "-d", "../data/unzipped/"])


def parsing_dependencies(conclude):
    sp = spacy.load('ja_ginza')
    doc = sp(conclude)
    word_dic = {}
    for token in doc:
        word_dic[token] = 0
    for token in doc:
        word_dic[token.head] += 1
    max_k = max(word_dic, key=word_dic.get)
    return max_k


def choice_color_direct(input_word, color_patterns, colors):
    impressions = color_patterns['name']
    indexes = [i for i, name in enumerate(impressions) if name == input_word]
    selected_impressions = color_patterns.iloc[indexes]
    selected_names = list(selected_impressions['name'])

    result_colors = get_color_images(
        colors, selected_impressions, indexes, selected_names)

    return result_colors[0]


def get_color_images(colors, impressions, indexes, names):
    result_colors = []
    for i, _ in enumerate(indexes):
        used_colors = list(impressions.iloc[i][impressions.iloc[i] == 1].index)
        result_a_color = []
        for j, color in enumerate(used_colors):
            name_rgb = colors[colors['Hue/Tone'] == color]
            rgb = [int(name_rgb['R']), int(name_rgb['G']), int(name_rgb['B'])]
            result_a_color.append(rgb)
        result_colors.append(result_a_color)

    return result_colors


def extract_origin(word):
    m = MeCab.Tagger('-Ochasen')
    node = m.parseToNode(word)
    while node:
        origin = node.feature.split(",")[6]
        if origin == "*":
            node = node.next
        else:
            return origin


def create_similar(word):
    word_origin = extract_origin(word)
    impressions = list(dict.fromkeys(list(color_patterns['name'])))
    similar_dict = {}
    for impression in impressions:
        impression_origin = extract_origin(impression)
        try:
            similar_dict[impression] = model.wv.similarity(word_origin,
                                                           impression_origin)
        except KeyError:
            pass
    return sorted(similar_dict.items(), key=lambda x: x[1], reverse=True)


def pptx_text(fname):
    f = open('sample.txt', 'w')
    # fname='sampleFile.pptx'
    txt = []
    prs = pptx.Presentation(fname)

    for i, sld in enumerate(prs.slides, start=1):

        for shp in sld.shapes:

            if shp.has_text_frame:
                shp.text = sub(shp.text)
                txt.append(shp.text)
                f.write(shp.text)

            if shp.has_table:
                tbl = shp.table
                row_count = len(tbl.rows)
                col_count = len(tbl.columns)
                for r in range(0, row_count):
                    text = ''
                    for c in range(0, col_count):
                        cell = tbl.cell(r, c)
                        paragraphs = cell.text_frame.paragraphs
                        for paragraph in paragraphs:
                            for run in paragraph.runs:
                                text += run.text
                            text += ', '
                    print(text)
                    f.write(text)
    f.close()
    sentences = [s for s in txt if s != ' ']
    sentences = [s for s in sentences if s != '']
    return sentences


def keyword():
    # ファイルパス
    file_path = "sample.txt"

    # 日本語テキストの読み込み
    f = open(file_path, "r", encoding="utf-8")
    text = f.read()
    f.close
    words = []
    # 形態素解析器で日本語処理
    t = Tokenizer()
    tokenize_text = t.tokenize(text)

    # Frequency生成＝複合語抽出処理（ディクショナリとリストの両方可)
    frequency = termextract.janome.cmp_noun_dict(tokenize_text)

    # FrequencyからLRを生成する
    lr = termextract.core.score_lr(
        frequency,
        ignore_words=termextract.mecab.IGNORE_WORDS,
        lr_mode=1, average_rate=1)

    # FrequencyとLRを組み合わせFLRの重要度を出す
    term_imp = termextract.core.term_importance(frequency, lr)

    # collectionsを使って重要度が高い順に表示
    data_collection = collections.Counter(term_imp)
    for cmp_noun, value in data_collection.most_common():
        words.append(termextract.core.modify_agglutinative_lang(cmp_noun))
    return words


def morpheme(sentences):
    # Tokenizerインスタンスの生成
    t = Tokenizer()
    # print(text)
    # テキストを引数として、形態素解析の結果、名詞・形容詞(原形)のみを配列で抽出する関数を定義

    def extract_words(text):
        tokens = t.tokenize(text)
        return [token.base_form for token in tokens
                if token.part_of_speech.split(',')[0] in ['名詞', '形容詞']]

    # それぞれの文章を単語リストに変換
    word_list = [extract_words(s) for s in sentences]
    words = []
    for word in word_list:
        # print(word)
        words.append(word)
    return words


def tf_idf():  # filename...モデルの保存ファイルが含まれるzipファイル
    for filename in glob.glob('*.zip'):
        # zipファイル解凍
        zip_f = zipfile.ZipFile(filename)
        zip_f.extractall('./data')
        zip_f.close()
    # データフレームに表現
    X = joblib.load('./data/X.sav')
    # データフレームに表現
    values = X.toarray()
    feature_names = joblib.load('./data/finalized_model.sav')
    df = pd.DataFrame(values, columns=feature_names,
                      index=["サンプル", "心理学", "情報科学", "経済", "社会学", "医学", "化学", "数学", "生物"])
    # 行列を転置
    df_0 = df[0:1].T
    # 値で降順ソート
    df_0 = df_0.sort_values(by="サンプル", ascending=False)
    return list(df_0.head(10).index)


if __name__ == '__main__':
    pass
